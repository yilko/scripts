import configparser;
import os;
import random;
from datetime import datetime;
from docx import Document;
from openpyxl import Workbook, load_workbook;
from pptx import Presentation;
from PIL import Image, ImageFont, ImageDraw;
from reportlab.pdfgen import canvas;
from PyPDF2 import PdfFileReader, PdfFileWriter;
from config_log import ConfigLog;
from save_files import SaveFiles;

'''实现创建及更新指定路径下的文件内容'''

log = ConfigLog().config_log();


class NewFiles:

    def __init__(self):
        # 初始化对象
        config = configparser.RawConfigParser();
        # 读取配置文件
        config.read("file_config.ini", encoding="utf-8-sig");
        self.is_create = config.getboolean("create_file", "is_create");
        self.name = config.get("create_file", "file_name");
        self.num = config.getint("create_file", "file_nums");
        self.content = config.get("create_file", "file_content");
        self.file_type = config.get("create_file", "file_type");
        self.folder = config.get("create_file", "folder_name");
        self.font_lib = config.get("create_file", "font_lib");
        self.is_update = config.getboolean("original_file", "is_update");

    # 返回新建文件夹下的文件列表
    def __new_folder_files(self) -> list:
        file_ls = [];
        # 列出当前目录下所有文件
        files = os.listdir(self.folder);
        log.debug(f"{self.folder}目录下有文件==={files}");
        for file in files:
            suffix = file.split(".")[-1];
            if suffix in ["docx", "xlsx", "pptx", "txt", "pdf"]:
                file_ls.append(file);
        log.debug(f"属于文件类型的有{len(file_ls)}个,分别是{file_ls}");
        return file_ls;

    # 获取创建文件的类型，并随机返回一个
    @staticmethod
    def __get_choice(file_type: str) -> str:
        if file_type.endswith("-"):
            file_type = file_type[:-1];
        file_ls = file_type.split("-");
        suffix = random.choice(file_ls);
        return suffix;

    # 几种文件类型中相同的部分抽离
    @staticmethod
    def __abstract_file(count: int, name: str, suffix: str, check_ls: list) -> tuple:
        file_name = f"{name}_{count}.{suffix}";
        time_now = datetime.now();
        timestamp = int(time_now.timestamp() * 1000);
        if file_name in check_ls:
            log.debug(f"存在冲突名字的文件==={file_name},已加时间戳重命名");
            file_name = f"{name}_{count}_{timestamp}.{suffix}";
        return file_name, time_now;

    # 创建docx格式
    def write_docx(self, mode: str, file_name: str, time_now: datetime):
        if mode == "create":
            # 不写参数默认新创建一个docx，如果同名在后面追加内容
            word = Document();
        elif mode == "update":
            word = Document(f"{self.folder}/{file_name}");
        else:
            raise Exception(f"既不是创建也不是修改，不支持的参数===={mode}");
        word.add_paragraph(self.content);
        word.add_paragraph(str(time_now));
        word.save(f"{self.folder}/{file_name}");
        log.info(f"{file_name}文件已{mode}完毕");

    # 创建xlsx格式
    def write_xlsx(self, mode: str, file_name: str, time_now: datetime):
        if mode == "create":
            wb = Workbook();
        elif mode == "update":
            wb = load_workbook(f"{self.folder}/{file_name}");
        else:
            raise Exception(f"既不是创建也不是修改，不支持的参数===={mode}");
        date = time_now.strftime("%F");
        ws = wb.create_sheet(date, 0);
        ws.cell(1, 1, self.content);
        ws.cell(2, 1, str(time_now));
        wb.save(f"{self.folder}/{file_name}");
        log.info(f"{file_name}文件已{mode}完毕");

    # 创建pptx格式
    def write_pptx(self, mode: str, file_name: str, time_now: datetime):
        if mode == "create":
            # 不写参数默认新创建一个ppt，如果同名在后面追加内容
            pre = Presentation();
        elif mode == "update":
            pre = Presentation(f"{self.folder}/{file_name}");
        else:
            raise Exception(f"既不是创建也不是修改，不支持的参数===={mode}");
        # 设置ppt的母版板式,加到幻灯片
        slide_template = pre.slide_layouts[0];
        slide = pre.slides.add_slide(slide_template);
        # 获取幻灯片的方框id(该段内容调试用)
        # for shape in slide.placeholders:
        #     phf = shape.placeholder_format;
        #     print(f"{phf.idx}--{shape.name}--{phf.type}");
        slide.placeholders[0].text = str(time_now);
        slide.placeholders[1].text = self.content;
        pre.save(f"{self.folder}/{file_name}");
        log.info(f"{file_name}文件已{mode}完毕");

    # 创建txt格式
    def write_txt(self, mode: str, file_name: str, time_now: datetime):
        with open(f"{self.folder}/{file_name}", "a", encoding="utf-8-sig") as f:
            f.write(f"{self.content}\n");
            f.write(f"{time_now}\n");
        log.info(f"{file_name}文件已{mode}完毕");

    # 创建图片
    def __txt_to_image(self, content: str, time_now: datetime) -> str:
        # 设置背景色白色，字体黑色
        background_color = "#FFFFFF";
        font_color = "#000000";
        img_size = (500, 200);
        # 创建一个图片对象
        image = Image.new("RGB", img_size, background_color);
        # 设置字体的格式大小(要simsun.ttc这种字体才支持中文)
        font = ImageFont.truetype(self.font_lib, int(img_size[1] * 0.1));
        # 创建图片中的画笔
        pen = ImageDraw.Draw(image);
        text = f"{content}\n{str(time_now)}";
        # draw.text(position, string, options)
        # 往图片中输入内容
        pen.text((0, 0), text, font=font, fill=font_color);
        img_name = "temp_photo.jpg";
        image.save(f"{self.folder}/{img_name}");
        return img_name;

    # 图片转成pdf(生成临时pdf或正式用的odf)
    def __img_to_pdf(self, pdf_name: str, time_now: datetime):
        img = self.__txt_to_image(self.content, time_now);
        # 打开图片并获取大小
        size = Image.open(f"{self.folder}/{img}").size;
        # 创建一个画布(pdf对象)，如果同名会覆盖原有内容
        can = canvas.Canvas(pdf_name, pagesize=size);
        # 往pdf中写入图片,从(0,0)开始(左下角)，右边两个参数分别是图片的宽和高
        can.drawImage(f"{self.folder}/{img}", 0, 0, size[0], size[1]);
        # 关闭并保存pdf
        can.showPage();
        can.save();

    # 创建pdf格式
    def write_pdf(self, mode: str, file_name: str, time_now: datetime):
        if mode == "create":
            pdf_name = f"{self.folder}/{file_name}";
            self.__img_to_pdf(pdf_name, time_now);
            log.info(f"{file_name}文件已{mode}完毕");
        # 临时pdf = 原pdf内容 + 时间戳pdf
        # 更新后的pdf = 临时pdf改名为原pdf名字
        elif mode == "update":
            original_pdf = f"{self.folder}/{file_name}";
            timestamp_pdf = f"{self.folder}/timestamp.pdf";
            self.__img_to_pdf(timestamp_pdf, time_now);
            pdf_writer = PdfFileWriter();
            # 读取原pdf和时间戳pdf的内容并写入临时pdf
            try:
                original_open = open(original_pdf, "rb");
                timestamp_open = open(timestamp_pdf, "rb");
                original_reader = PdfFileReader(original_open, strict=False);
                timestamp_reader = PdfFileReader(timestamp_open, strict=False);
            except Exception as e:
                log.error(e);
                log.error(f"该文件不支持打开编辑内容===={original_pdf}");
            else:
                for reader in [original_reader, timestamp_reader]:
                    nums = reader.getNumPages();
                    for num in range(nums):
                        page = reader.getPage(num);
                        pdf_writer.addPage(page);
                # 保存临时pdf文件
                with open(f"{self.folder}/temp.pdf", "wb") as f:
                    pdf_writer.write(f);
                original_open.close();
                timestamp_open.close();
            try:
                # 删除原pdf和时间戳pdf文件
                for pdf in [original_pdf, timestamp_pdf]:
                    if os.path.exists(pdf):
                        os.remove(pdf);
                # 临时pdf命名为原pdf内容
                os.rename(f"{self.folder}/temp.pdf", original_pdf);
                log.info(f"{file_name}文件已{mode}完毕");
            except Exception as e:
                log.error(e);
        else:
            raise Exception(f"既不是创建也不是修改，不支持的参数===={mode}");

    # 开始随机创建文件
    def create_files(self):
        if self.is_create:
            if not os.path.exists(self.folder):
                os.mkdir(self.folder);
            # 返回新建文件夹的文件列表
            files_ls = self.__new_folder_files();
            for count in range(self.num):
                # 随机生成后缀格式
                suffix = NewFiles.__get_choice(self.file_type);
                # 检查重名并返回文件名称，时间戳
                file_name, time_now = NewFiles.__abstract_file(count, self.name, suffix, files_ls);
                func_name = f"write_{suffix}";
                # 反射调用对应函数
                if hasattr(self, func_name):
                    func = getattr(self, func_name);
                    func("create", file_name, time_now);
                else:
                    log.error(f"不支持创建的格式===={suffix},无法调用对应方法创建文件===={func_name}");
            else:
                # 删除临时img图片
                img_path = f"{self.folder}/temp_photo.jpg";
                if os.path.exists(img_path):
                    os.remove(img_path);
                log.debug(f"-----------------删除临时图片===={img_path}-----------------");
                log.info(f"--------------总共给{self.num}个文件进行创建----------------");

        else:
            log.info("不需要新建文件");

    # 原有文件进行内容更新
    def update_file(self):
        if self.is_update:
            if os.path.exists(self.folder):
                file_ls = self.__new_folder_files();
                if file_ls:
                    for file in file_ls:
                        suffix = file.split(".")[-1];
                        func_name = f"write_{suffix}";
                        if hasattr(self, func_name):
                            func = getattr(self, func_name);
                            func("update", file, datetime.now());
                    else:
                        # 删除临时img图片
                        img_path = f"{self.folder}/temp_photo.jpg";
                        if os.path.exists(img_path):
                            os.remove(img_path);
                        log.debug(f"-----------------删除临时图片===={img_path}-----------------");
                        log.debug(f"文件更新列表===={file_ls}");
                        log.info(f"--------------总共给{len(file_ls)}个文件进行更新---------------");
                else:
                    log.info(f"文件夹下没有文件类型(docx,xlsx,pptx,txt,pdf)，不需要更新内容,列表为空===={file_ls}");
            else:
                log.info(f"当前路径下没有创建{self.folder}文件夹，没有文件可更新内容");
        else:
            log.info("不需要更新原有文件内容");


if __name__ == '__main__':
    try:
        nf = NewFiles();
        nf.update_file();
        nf.create_files();
        sf = SaveFiles();
        sf.save_md5();
        os.startfile(nf.folder);
        os.system("pause");
    except Exception as e:
        log.error(e);
        os.system("pause");
